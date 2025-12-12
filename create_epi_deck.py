from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(title, content_items, layout_index=1):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add Content (Bullet points)
        if content_items:
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()  # Clear existing empty paragraph
            
            for item in content_items:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.font.size = Pt(18)

    # Helper for Title Slide
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    # 1. Title Slide
    add_title_slide(
        "EPI: Evidence Packaged Infrastructure",
        "\"Git for AI Runs\" - Turn opaque AI executions into cryptographically-verified, shareable evidence packages.\n\nPre-Seed Pitch Deck"
    )

    # 2. The Problem
    add_slide(
        "The Problem: The AI \"Black Box\" Crisis",
        [
            "Reproducibility Crisis: 70% of AI research papers cannot be reproduced.",
            "Production Failures: \"It worked on my machine\" is rampant; debugging production AI is expensive.",
            "Compliance Gaps: EU AI Act & FDA require audit trails that don't exist today.",
            "Trust Deficit: AI decisions are opaque with no chain of custody."
        ]
    )

    # 3. The Solution
    add_slide(
        "The Solution: EPI (Evidence Packaged Infrastructure)",
        [
            "The \"Black Box Recorder\" for AI systems.",
            "Captures a complete timeline of every execution (LLM calls, inputs, outputs, env).",
            "Creates a single, portable .epi file.",
            "Cryptographically signed and independently verifiable.",
            "Includes an embedded, serverless HTML viewer."
        ]
    )

    # 4. The PDF Analogy
    add_slide(
        "The PDF Analogy",
        [
            "Just as PDF became the standard for Documents:",
            "  - Self-contained & Platform-independent",
            "  - Preserves formatting vs. Preserves execution context",
            "  - Signed & Verifiable",
            "",
            "EPI is the standard for AI Workflows.",
            "\"Every AI decision comes with an .epi file.\""
        ]
    )

    # 5. Key Product Features
    add_slide(
        "Product Features",
        [
            "Zero-Config Recording: Auto-patching of OpenAI/LangChain SDKs.",
            "Universal Format: Works on Windows, macOS, Linux.",
            "Security: Automatic redaction of API keys (15+ patterns).",
            "Cryptography: Ed25519 signatures & SHA-256 hashing.",
            "Interactive Viewer: Beautiful timeline UI within the file itself."
        ]
    )

    # 6. Target Customers
    add_slide(
        "Target Customers",
        [
            "AI Researchers: Verify claims and attach .epi files to papers.",
            "AI Startups: Debug production failures 10x faster.",
            "Enterprises: Compliance with EU AI Act and internal audits.",
            "Healthcare/Finance: FDA/SEC audit trails and liability defense."
        ]
    )

    # 7. Competitive Positioning
    add_slide(
        "Competitive Landscape",
        [
            "vs. LangSmith/Helicone: They are monitoring tools (cloud, logs). EPI is evidence (local, signed file).",
            "vs. Git: Git tracks code. EPI tracks execution.",
            "vs. MLflow: EPI is for compliance and evidence, not just experiment tracking.",
            "",
            "EPI is complementary: Use LangSmith for monitoring, EPI for the permanent record."
        ]
    )

    # 8. Business Model (Open Core)
    add_slide(
        "Business Model: Open Core",
        [
            "Free / Open Source:",
            "  - CLI Tool, Python API, Local Recording.",
            "  - Community growth & standard creation.",
            "",
            "Commercial (Future SaaS/Enterprise):",
            "  - Cloud Storage & Collaboration.",
            "  - Advanced Analytics & Team Management.",
            "  - Enterprise On-Premise & SSO.",
            "  - Compliance Reporting Automation."
        ]
    )

    # 9. Vision & Roadmap
    add_slide(
        "Vision & Roadmap",
        [
            "Phase 1 (Now): Open Source Traction (PyPI, Hacker News).",
            "Phase 2: Research Adoption (Partner with Labs).",
            "Phase 3: Enterprise POCs (Healthcare/Finance pilots).",
            "Phase 4: The Standard (`.epi` everywhere).",
            "",
            "Vision: Total transparency for AI. Solving the reproducibility crisis."
        ]
    )

    # 10. Contact / Closing
    add_title_slide(
        "Investing in Trust",
        "EPI Labs\ncontact@epilabs.org\nhttps://github.com/mohdibrahimaiml/EPI-OFFICIAL"
    )

    output_file = "EPI_Pitch_Deck.pptx"
    prs.save(output_file)
    print(f"Successfully created {output_file}")

if __name__ == "__main__":
    create_presentation()
