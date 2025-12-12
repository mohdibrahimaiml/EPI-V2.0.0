from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper to add a title slide
    def add_title_slide(title, subtitle, tagline, details):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        # Subtitle and Tagline
        subtitle_shape = slide.placeholders[1]
        tf = subtitle_shape.text_frame
        tf.clear()
        
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.bold = True
        p.font.size = Pt(28)
        
        p = tf.add_paragraph()
        p.text = tagline
        p.font.size = Pt(24)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "\n" + details
        p.font.size = Pt(18)

    # Helper to add a content slide with text
    def add_content_slide(title, content_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        
        for item in content_list:
            p = tf.add_paragraph()
            if isinstance(item, tuple):
                text, level = item
                p.text = text
                p.level = level
            else:
                p.text = item
                p.level = 0
            p.font.size = Pt(14)

    # Helper to create a table slide
    def add_table_slide(title, table_data, summary_text=None):
        slide_layout = prs.slide_layouts[5] # Title Only
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        # Table calculation
        rows = len(table_data)
        cols = len(table_data[0])
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(0.8) # arbitrary
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Determine column widths
        # Just simple distribution for now
        for i in range(cols):
            table.columns[i].width = int(width / cols)

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(table_data[r][c])
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                if r == 0:
                    cell.text_frame.paragraphs[0].font.bold = True

        if summary_text:
            textbox = slide.shapes.add_textbox(Inches(1), top + Inches(rows * 0.4) + Inches(0.5), Inches(8), Inches(2))
            tf = textbox.text_frame
            tf.word_wrap = True
            for line in summary_text:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)

    # --- SLIDES GENERATION ---

    # SLIDE 1: TITLE SLIDE
    add_title_slide(
        "Evidence Packaged Infrastructure (EPI)",
        "The Infrastructure Standard for AI Reproducibility",
        "Record -> Sign -> Replay -> Prove",
        "Founder: Mohd Ibrahim Afridi\nEPI Labs | Hyderabad, TS\nFounded: 2024\nRaising: $250K USD (Preseed)"
    )

    # SLIDE 2: THE PROBLEM
    add_content_slide(
        "The Problem: Regulators Can't Verify Why AI Models Make Decisions",
        [
            "The Core Problem:",
            ("AI systems in regulated industries (banking, pharma) are opaque 'black boxes'.", 1),
            ("Regulators (RBI, FDA, EU AI Act) cannot verify reproducibility.", 1),
            ("Logs exist, but cannot prove exact computation was performed.", 1),
            ("Non-compliance costs: FDA device rejection, RBI audit failure, GDPR fines.", 1),
            "",
            "Regulatory Pressure (Real & Immediate):",
            ("RBI Phase 1 (NOW): Audit trails required.", 1),
            ("FDA (May 2025): AI/ML device guidance issued.", 1),
            ("EU AI Act (2025-26): Explainability proof required.", 1),
            "",
            "Customer Evidence:",
            ("Fintech CEO: 'We can't prove to regulators why our model made this credit decision.'", 1),
            ("Biotech CRO: 'FDA is asking us to prove model reproducibility.'", 1)
        ]
    )

    # SLIDE 3: THE SOLUTION
    add_content_slide(
        "The Solution: EPI - Infrastructure for Reproducibility Proof",
        [
            "What is EPI?",
            ("A lightweight infrastructure primitive that records, signs, and replays AI execution.", 1),
            "",
            "Three Core Capabilities:",
            ("1. Record: Captures complete execution (inputs, weights, outputs) with minimal code.", 1),
            ("2. Sign: Cryptographically signs execution proof (Ed25519) - tamper-proof.", 1),
            ("3. Replay: Recreates exact model execution from signature (100% reproducibility).", 1),
            "",
            "Why This Matters:",
            ("Only vendor combining Record + Sign + Replay.", 1),
            ("Portable .epi format accepted by regulators.", 1),
            ("Open-source foundation (trust, transparency).", 1),
            "",
            "Product Status:",
            ("MVP on PyPI, Replay demo working, 60+ test cases.", 1)
        ]
    )

    # SLIDE 4: MARKET SIZE
    market_data = [
        ["Segment", "Market (2025)", "EPI SAM"],
        ["Enterprise AI Gov", "$2.2B", "$330M"],
        ["AI Gov (Narrow)", "$300M+", "$60M"],
        ["MLOps", "$2.3B", "$276M"],
        ["Audit Software", "$10B", "$650M"],
        ["Pharma (FDA)", "$3B", "$300M"],
        ["Others (RBI, EU)", "$500B+", "$150M+"]
    ]
    add_table_slide(
        "Market Size: $1.8B Global TAM (2027-2030)",
        market_data,
        [
            "Preseed TAM (Available NOW): $58-95M",
            "- Fintech (RBI Phase 1): $3-5M",
            "- Pharma (FDA Guidance): $20-30M",
            "- Enterprise (GDPR/CCPA): $30-50M",
            "",
            "Key Insight for Rainmatter:",
            "Preseed isn't chasing the $1.8B standard play yet. We capture the $34-55M from early adopters feeling pressure NOW. Series A opens when Phase 2 begins (FY26-27)."
        ]
    )

    # SLIDE 5: MARKET TIMING
    add_content_slide(
        "Market Timing: Regulatory Mandates Create 3-5x Expansion",
        [
            "Phase 1 (2025 - NOW): Compliance Awakening",
            ("RBI Phase 1 live, FDA guidance issued.", 1),
            ("Market: $34-55M. Target: Land 3-5 design partners.", 1),
            "",
            "Phase 2 (FY26-27): Mandatory Audit Frameworks",
            ("RBI Phase 2 implies mandatory audits. FDA enforcement begins.", 1),
            ("Market: $150-250M (3-5x expansion).", 1),
            "",
            "Phase 3 (2027-28): Mandatory Compliance + Penalties",
            ("RBI Phase 3: All AI models need proof. High-risk FDA devices need proof.", 1),
            ("Market: $400-600M.", 1),
            "",
            "Why Rainmatter?",
            ("Fintech portfolio benefits from Phase 1 compliance.", 1),
            ("Patient capital for infrastructure play.", 1)
        ]
    )

    # SLIDE 6: COMPETITION
    comp_data = [
        ["Capability", "EPI", "Databricks", "AWS", "MLflow"],
        ["Record Exec", "YES", "Add-on", "Add-on", "NO"],
        ["Crypto Sign", "YES", "NO", "NO", "NO"],
        ["Replay Exact", "YES", "NO", "NO", "NO"],
        ["Audit Proof", "YES", "Partial", "Partial", "NO"],
        ["Portable", "YES", "NO", "NO", "Partial"]
    ]
    add_table_slide(
        "Why EPI Wins (vs Databricks, AWS, MLflow)",
        comp_data,
        [
            "Why Large Vendors Won't Win:",
            "- Databricks/AWS: Reproducibility is a secondary feature. Compliance is an add-on.",
            "- MLflow: Logs data but cannot sign/replay (not audit-grade).",
            "",
            "Why EPI Has a 2-3 Year Window:",
            "- First-mover in standardized reproducibility proof.",
            "- Format moat (.epi) vs Platform lock-in.",
            "- We own the regulation conversation (RBI/FDA)."
        ]
    )

    # SLIDE 7: TRACTION
    add_content_slide(
        "Traction & Validation: MVP Built, Market Validated",
        [
            "Product Traction:",
            ("epi-recorder on PyPI (Production Ready).", 1),
            ("GitHub: 60+ automated tests, Replay demo working.", 1),
            "",
            "Customer Validation (Preseed Convos):",
            ("Fintech CEO: 'RBI Phase 1 is asking us for audit trails. EPI does exactly what we need.'", 1),
            ("Biotech CRO: 'Your replay demo solves our FDA problem.'", 1),
            ("", 1),
            "Growth Trajectory:",
            ("Q4 2025: MVP Done, 3 Design Partner Convos.", 1),
            ("Q1 2026: Close 1st Partner ($100K).", 1),
            ("Q2-Q4 2026: Scale to $900K-$1.5M ARR.", 1)
        ]
    )

    # SLIDE 8: BUSINESS MODEL
    add_content_slide(
        "Business Model: How We'll Acquire Customers",
        [
            "Preseed Strategy: Design Partner Playbook",
            ("Target: 3-5 partners (Fintech, Pharma). Price: $75-150K/year.", 1),
            ("Channel: Founder-led sales. Timeline: 3-9 months.", 1),
            "",
            "Revenue Model:",
            ("SaaS Subscription (80%): Tier 1 ($75K), Tier 2 ($150K), Tier 3 ($300K+).", 1),
            ("Managed Services (15%): Audit setup & validation ($50-100K).", 1),
            ("Licensing (5%): Cloud hosting & verification.", 1),
            "",
            "Unit Economics (Exceptional):",
            ("CAC Payback: 3-4 months (vs 12mo standard).", 1),
            ("Retention: >95% (Compliance makes it sticky).", 1),
            ("NRR: 150% (Expansion revenue).", 1)
        ]
    )

    # SLIDE 9: FINANCIALS & ASK
    fin_data = [
        ["Metric", "Preseed (25-26)", "Series A (26-27)", "Series B"],
        ["ARR", "$0.5-1.5M", "$2-3M", "$10-20M"],
        ["Customers", "4-8", "15-20", "50-100"],
        ["ARPU", "$100-150K", "$120-150K", "$150-200K"]
    ]
    add_table_slide(
        "Financial Forecast & The Ask",
        fin_data,
        [
            "Raising: $250K USD (Preseed) | Valuation: $3-5M",
            "Use of Funds:",
            "- 40% Engineering (Product hardening, Signing)",
            "- 30% Sales & Marketing (Design partners)",
            "- 16% Ops, 14% Runway",
            "",
            "Investor Returns (Rainmatter Stake):",
            "- Base Case Exit ($1B): 240-360x return ($60-90M).",
            "- Strategic Value: Reproducibility compliance for Zerodha ecosystem."
        ]
    )

    output_file = "EPI_Rainmatter_Pitch_Deck.pptx"
    prs.save(output_file)
    print(f"Successfully created {output_file}")

if __name__ == "__main__":
    create_presentation()
